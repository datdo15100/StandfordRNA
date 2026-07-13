"""Generate the ready-to-present PowerPoint deck from repository results."""
from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


HERE = Path(__file__).resolve().parent
CHARTS = HERE / "charts"
CHARTS.mkdir(exist_ok=True)

NAVY = "14213D"
BLUE = "2F80ED"
TEAL = "13B8A6"
ORANGE = "F2994A"
RED = "D64550"
INK = "172033"
MUTED = "5E6B7A"
PALE = "F4F7FB"
WHITE = "FFFFFF"


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def style_chart() -> None:
    plt.rcParams.update(
        {
            "font.family": "DejaVu Sans",
            "font.size": 11,
            "axes.spines.top": False,
            "axes.spines.right": False,
            "axes.edgecolor": "#AAB4C3",
            "axes.labelcolor": f"#{INK}",
            "xtick.color": f"#{MUTED}",
            "ytick.color": f"#{MUTED}",
        }
    )


def save_charts() -> dict[str, Path]:
    style_chart()
    charts: dict[str, Path] = {}

    stages = ["Dummy", "TBM + refine", "+ de novo", "+ composite", "Top-1 reproduced"]
    scores = [0.069, 0.161, 0.212, 0.3072, 0.2973]
    colors = [f"#{MUTED}", f"#{BLUE}", f"#{TEAL}", f"#{ORANGE}", f"#{NAVY}"]
    fig, ax = plt.subplots(figsize=(9, 4.5))
    bars = ax.bar(stages, scores, color=colors, width=0.7)
    ax.set_ylabel("Mean best-of-5 TM")
    ax.set_ylim(0, 0.35)
    ax.grid(axis="y", alpha=0.18)
    ax.bar_label(bars, labels=[f"{v:.3f}" for v in scores], padding=3, fontsize=11)
    ax.tick_params(axis="x", rotation=12)
    fig.tight_layout()
    charts["evolution"] = CHARTS / "pipeline_evolution.png"
    fig.savefig(charts["evolution"], dpi=180, transparent=False, facecolor="white")
    plt.close(fig)

    regimes = ["Temporal-safe", "Ignore cutoff\n(exclude native)", "Native/oracle\nallowed"]
    values = [0.1612, 0.6388, 0.9566]
    fig, ax = plt.subplots(figsize=(7.5, 4.2))
    bars = ax.bar(regimes, values, color=[f"#{TEAL}", f"#{ORANGE}", f"#{RED}"])
    ax.set_ylabel("Mean best-of-5 TM")
    ax.set_ylim(0, 1.05)
    ax.grid(axis="y", alpha=0.18)
    ax.bar_label(bars, labels=[f"{v:.3f}" for v in values], padding=3)
    fig.tight_layout()
    charts["leakage"] = CHARTS / "leakage_diagnostic.png"
    fig.savefig(charts["leakage"], dpi=180, facecolor="white")
    plt.close(fig)

    targets = ["R1107", "R1108", "R1116", "R1117v2", "R1126", "R1128", "R1136", "R1138", "R1149", "R1156", "R1189", "R1190"]
    delta = [0.0545, 0.1765, -0.0084, 0.3200, 0.0446, 0.0465, 0.0721, 0.0510, 0.1599, 0.1064, 0.0569, 0.0664]
    fig, ax = plt.subplots(figsize=(9, 4.5))
    y = np.arange(len(targets))
    ax.barh(y, delta, color=[f"#{TEAL}" if v >= 0 else f"#{RED}" for v in delta])
    ax.set_yticks(y, targets)
    ax.invert_yaxis()
    ax.axvline(0, color=f"#{INK}", linewidth=0.8)
    ax.set_xlabel("ΔTM from composite search")
    ax.grid(axis="x", alpha=0.18)
    fig.tight_layout()
    charts["composite"] = CHARTS / "composite_delta.png"
    fig.savefig(charts["composite"], dpi=180, facecolor="white")
    plt.close(fig)

    methods = ["None", "Rule", "Gradient v1"]
    tm = [0.3092, 0.3098, 0.3072]
    clash = [0.1634, 0.0991, 0.0935]
    kink = [0.0536, 0.0944, 0.1025]
    fig, axes = plt.subplots(1, 3, figsize=(10.5, 3.8))
    for ax, vals, title, good in zip(
        axes,
        [tm, clash, kink],
        ["TM (higher better)", "Clashes/res (lower)", "Sharp kinks (lower)"],
        [TEAL, TEAL, RED],
    ):
        bars = ax.bar(methods, vals, color=[f"#{MUTED}", f"#{BLUE}", f"#{good}"])
        ax.set_title(title, fontsize=11)
        ax.bar_label(bars, labels=[f"{v:.3f}" for v in vals], padding=2, fontsize=9)
        ax.grid(axis="y", alpha=0.15)
        ax.tick_params(axis="x", rotation=20, labelsize=9)
        ax.set_ylim(0, max(vals) * 1.25)
    fig.tight_layout()
    charts["refinement"] = CHARTS / "refinement_truthfulness.png"
    fig.savefig(charts["refinement"], dpi=180, facecolor="white")
    plt.close(fig)
    return charts


class Deck:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def _background(self, slide, color: str = WHITE) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb(color)

    def _footer(self, slide, number: int, dark: bool = False) -> None:
        color = "C8D2E0" if dark else MUTED
        box = slide.shapes.add_textbox(Inches(0.55), Inches(7.12), Inches(12.2), Inches(0.2))
        p = box.text_frame.paragraphs[0]
        p.text = f"Stanford RNA 3D Folding thesis update  •  14 Jul 2026                                      {number}"
        p.font.size = Pt(8)
        p.font.color.rgb = rgb(color)

    def _title(self, slide, title: str, subtitle: str | None = None) -> None:
        box = slide.shapes.add_textbox(Inches(0.65), Inches(0.35), Inches(12.0), Inches(0.75))
        p = box.text_frame.paragraphs[0]
        p.text = title
        p.font.name = "Aptos Display"
        p.font.bold = True
        p.font.size = Pt(28)
        p.font.color.rgb = rgb(NAVY)
        if subtitle:
            sb = slide.shapes.add_textbox(Inches(0.68), Inches(1.04), Inches(11.8), Inches(0.4))
            sp = sb.text_frame.paragraphs[0]
            sp.text = subtitle
            sp.font.size = Pt(12)
            sp.font.color.rgb = rgb(MUTED)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.68), Inches(1.35), Inches(1.25), Inches(0.06))
        line.fill.solid(); line.fill.fore_color.rgb = rgb(TEAL); line.line.fill.background()

    def _bullets(self, slide, items: list[str], x=0.8, y=1.65, w=11.7, h=4.9, size=22) -> None:
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.05)
        tf.clear()
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.name = "Aptos"
            p.font.size = Pt(size)
            p.font.color.rgb = rgb(INK)
            p.space_after = Pt(11)
            p.line_spacing = 1.06

    def _callout(self, slide, text: str, x: float, y: float, w: float, h: float, color=TEAL, size=17) -> None:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid(); shape.fill.fore_color.rgb = rgb(PALE)
        shape.line.color.rgb = rgb(color); shape.line.width = Pt(1.5)
        tf = shape.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER
        p.font.name = "Aptos"; p.font.size = Pt(size); p.font.bold = True; p.font.color.rgb = rgb(INK)

    def title_slide(self) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]); self._background(slide, NAVY)
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
        accent.fill.solid(); accent.fill.fore_color.rgb = rgb(TEAL); accent.line.fill.background()
        box = slide.shapes.add_textbox(Inches(0.85), Inches(1.25), Inches(11.3), Inches(2.3))
        tf = box.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.text = "GeoFuse-RNA"; p.font.name = "Aptos Display"; p.font.size = Pt(46); p.font.bold = True; p.font.color.rgb = rgb(WHITE)
        p2 = tf.add_paragraph(); p2.text = "Confidence-aware fusion and motif-conditioned\ngeometric refinement for RNA 3D prediction"; p2.font.size = Pt(24); p2.font.color.rgb = rgb("DCE6F3"); p2.space_before = Pt(16)
        tag = slide.shapes.add_textbox(Inches(0.9), Inches(5.4), Inches(6), Inches(0.8))
        tp = tag.text_frame.paragraphs[0]; tp.text = "Supervisor progress update\n14 July 2026"; tp.font.size = Pt(17); tp.font.color.rgb = rgb("AFC0D6")
        self._footer(slide, 1, dark=True)

    def bullet_slide(self, title: str, items: list[str], number: int, subtitle: str | None = None, size=22) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]); self._background(slide)
        self._title(slide, title, subtitle); self._bullets(slide, items, size=size); self._footer(slide, number)

    def chart_slide(self, title: str, chart: Path, number: int, takeaway: str, subtitle: str | None = None) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]); self._background(slide)
        self._title(slide, title, subtitle)
        slide.shapes.add_picture(str(chart), Inches(0.75), Inches(1.58), width=Inches(8.7), height=Inches(4.65))
        self._callout(slide, takeaway, 9.7, 2.0, 2.85, 3.2, color=ORANGE, size=17)
        self._footer(slide, number)

    def flow_slide(self, number: int) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]); self._background(slide)
        self._title(slide, "Implemented pipeline", "One code path for local evaluation and Kaggle inference")
        labels = ["Sequence + cutoff", "Search", "Temporal filter", "Transfer + gap fill", "Candidate pool", "Final five"]
        xs = [0.55, 2.65, 4.65, 6.68, 8.88, 11.0]
        for i, (label, x) in enumerate(zip(labels, xs)):
            color = TEAL if i in (1, 4) else BLUE
            self._callout(slide, label, x, 2.15, 1.75, 1.15, color=color, size=14)
            if i < len(labels) - 1:
                arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 1.72), Inches(2.52), Inches(0.43), Inches(0.35))
                arrow.fill.solid(); arrow.fill.fore_color.rgb = rgb(MUTED); arrow.line.fill.background()
        self._bullets(slide, ["MMseqs + exhaustive composite similarity", "Leakage guard before ranking", "TBM + de novo hedge + optional geometry refinement", "Quality/diversity selection produces five C1′ structures"], x=1.0, y=4.05, w=11.3, h=2.3, size=18)
        self._footer(slide, number)

    def geofuse_slide(self, number: int) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]); self._background(slide, NAVY)
        title = slide.shapes.add_textbox(Inches(0.65), Inches(0.4), Inches(12), Inches(0.7))
        p = title.text_frame.paragraphs[0]; p.text = "Proposed thesis extension: GeoFuse-RNA"; p.font.size = Pt(29); p.font.bold = True; p.font.color.rgb = rgb(WHITE)
        boxes = [
            ("TBM candidates", 0.8, 1.8, BLUE), ("Pretrained candidates", 0.8, 4.1, TEAL),
            ("Fold clustering", 3.45, 2.95, ORANGE), ("Residue/segment confidence", 5.85, 2.95, TEAL),
            ("Segment fusion", 8.65, 2.15, BLUE), ("Geometry v2 projection", 8.65, 4.0, ORANGE),
            ("Quality-diversity final five", 11.0, 3.05, TEAL),
        ]
        for label, x, y, color in boxes:
            w = 1.85 if x < 3 else (2.1 if x < 8 else 1.95)
            self._callout(slide, label, x, y, w, 1.05, color=color, size=13)
        arrows = [(2.65, 2.2, 3.45, 3.25), (2.65, 4.45, 3.45, 3.55), (5.55, 3.35, 5.85, 3.35), (7.95, 3.35, 8.65, 2.7), (7.95, 3.55, 8.65, 4.45), (10.6, 2.65, 11.0, 3.3), (10.6, 4.4, 11.0, 3.7)]
        for x1, y1, x2, y2 in arrows:
            line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
            line.line.color.rgb = rgb("AFC0D6"); line.line.width = Pt(1.5)
        claim = slide.shapes.add_textbox(Inches(1.1), Inches(6.25), Inches(11.1), Inches(0.55))
        cp = claim.text_frame.paragraphs[0]; cp.text = "Novelty is the adaptive integration layer — not another large RNA foundation model."; cp.alignment = PP_ALIGN.CENTER; cp.font.size = Pt(18); cp.font.bold = True; cp.font.color.rgb = rgb(WHITE)
        self._footer(slide, number, dark=True)

    def experiment_slide(self, number: int) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]); self._background(slide)
        self._title(slide, "Experiments that isolate the contribution")
        rows = [
            ("B0", "Current TBM + composite", "Strong reference baseline"),
            ("B1", "+ raw pretrained union", "Does oracle candidate TM rise?"),
            ("B2", "+ confidence-aware fusion", "Fusion vs whole-candidate selection"),
            ("B3", "+ geometry v2", "Repair without kink artifacts"),
            ("B4", "+ fold-aware final-five", "Quality/diversity and selection regret"),
        ]
        y = 1.72
        for i, (tag, method, question) in enumerate(rows):
            color = TEAL if i >= 2 else BLUE
            self._callout(slide, tag, 0.8, y, 0.75, 0.72, color=color, size=16)
            self._callout(slide, method, 1.75, y, 4.0, 0.72, color=color, size=15)
            self._callout(slide, question, 6.05, y, 6.15, 0.72, color=ORANGE, size=15)
            y += 0.94
        self._footer(slide, number)

    def save(self, path: Path) -> None:
        self.prs.save(path)


def main() -> None:
    charts = save_charts()
    deck = Deck()
    deck.title_slide()
    deck.bullet_slide("Executive summary", [
        "Reproducible temporal-safe RNA 3D pipeline built around the Stanford Kaggle challenge.",
        "Strongest local result: 0.3072 mean best-of-5 TM on 12 CASP15 targets.",
        "Main empirical win: +0.0955 TM from better template recall.",
        "Geometry v1 improves clashes/backbone spacing, but is TM-neutral and increases sharp kinks.",
        "Next contribution: segment-level TBM/pretrained fusion plus motif-conditioned geometry v2.",
    ], 2, size=21)
    deck.bullet_slide("Problem and benchmark", [
        "Input: RNA sequence, MSA, template structures and release-date metadata.",
        "Output: five C1′ coordinate structures per target, each with shape [L, 3].",
        "Competition score: best TM among five predictions, averaged over targets/references.",
        "TM rewards the global fold and can tolerate local geometric artifacts.",
        "Thesis evaluation therefore separates fold accuracy from structural validity.",
    ], 3)
    deck.bullet_slide("Data and EDA snapshot", [
        "Train v1: 844 / 137k residues; train v2: 5,135 / 3.68M residues.",
        "Local validation/test: 12 CASP15 targets, 2,515 residues total.",
        "Lengths span 30–720 nt (median 129.5 nt); R1138 is a 720-nt stress case.",
        "Template resource: 8,670 RNA CIF files, 56.89 GiB.",
        "Historical parse: 23,869 chains, 10.86M residues, 99.9% modelled C1′, zero errors.",
        "MSA and MSA_v2 add target-specific evolutionary context.",
    ], 4, size=20)
    deck.chart_slide("Leakage is the central evaluation risk", charts["leakage"], 5, "Ignoring time can inflate TM by +0.48 to +0.80.\n\nOracle results are diagnostics, not scientific performance.", "Same earlier TBM pipeline under three template-availability regimes")
    deck.flow_slide(6)
    deck.chart_slide("How performance evolved", charts["evolution"], 7, "Candidate generation dominates.\n\nThe largest jump came from search recall, not coordinate optimisation.")
    deck.bullet_slide("Bottleneck diagnosis: candidate recall", [
        "Only 5/12 targets originally had an MMseqs temporal-safe hit; 7/12 fell back to de novo.",
        "A faithful top-1 reproduction scored 0.2973 temporal-safe using exhaustive composite similarity.",
        "This diagnosed search—not refinement—as the main accuracy bottleneck.",
        "Composite search lifted our pipeline from 0.2117 to 0.3072.",
        "It improved 11/12 targets and beat reproduced top-1 on 9/12.",
    ], 8, size=21)
    deck.chart_slide("Composite-search ablation", charts["composite"], 9, "Mean gain: +0.0955 TM\n\nLargest: R1117v2 +0.320\n\nRuntime: ~8 s/target")
    deck.chart_slide("Geometry refinement v1: an honest negative result", charts["refinement"], 10, "−42% clashes\n−47% backbone deviation\n\nTM: −0.002\nKinks: ~2×", "Auxiliary metrics averaged over all five structures and 12 targets")
    deck.bullet_slide("What is already a contribution?", [
        "A reproducible temporal-safe benchmark with explicit leakage quantification.",
        "A search diagnosis and composite-recall improvement over the reproduced top-1 baseline.",
        "Adversarial refinement evaluation using an unoptimised sharp-kink metric.",
        "Reusable accuracy, geometry and diversity evaluation in one codebase.",
        "However, generic ‘TBM + pretrained + refinement’ remains too close to leaderboard practice.",
    ], 11, size=21)
    deck.geofuse_slide(12)
    deck.bullet_slide("Geometry v2: directly answer the v1 failure", [
        "Retain adjacent-distance, clash and size terms.",
        "Add context-conditioned angle/curvature distributions to prevent kink trading.",
        "Use signed pseudo-torsion only when the required atom representation is retained.",
        "Optimise in stages: stitch/fuse → local repair → weak global prior.",
        "Success gate: preserve TM and geometry gains without exceeding the no-refine kink rate.",
    ], 13, size=21)
    deck.bullet_slide("Confidence-aware segment fusion", [
        "Estimate source reliability at residue/segment level, not only one score per structure.",
        "TBM features: identity, coverage, completeness, gap mask, date and local agreement.",
        "Pretrained features: model confidence and agreement among predicted folds.",
        "Align and cluster candidates before fusion; never blend incompatible global folds.",
        "Use segment smoothing and seam penalties to avoid ‘Frankenstein’ switching.",
        "Train/evaluate on real held-out predictions, not synthetic corruptions alone.",
    ], 14, size=19)
    deck.experiment_slide(15)
    deck.bullet_slide("Current engineering status", [
        "WSL Conda env ‘rna-fold’ installed; RTX 3060 Ti CUDA and MMseqs2 verified.",
        "Four core numerical/temporal tests pass.",
        "Kaggle token and CLI work; the account currently has no competition submission.",
        "Raw data validated: 61 GiB with all CSV, MSA and 8,670 PDB CIF components present.",
        "WSL cap configured to 18 GB RAM + 8 GB swap; restart is required to apply it.",
        "CIF rebuild defaults to six workers to preserve Windows headroom.",
    ], 16, size=20)
    deck.bullet_slide("Compute strategy and laptop handoff", [
        "RTX 3060 Ti / 8 GB: artifact builds and DRfold2/RibonanzaNet-scale experiments.",
        "GTX 1650 / 4 GB laptop: code, unit tests, cached-candidate analysis, slides and thesis writing.",
        "Do not rebuild 57 GB CIFs or run AF3-style models on the laptop.",
        "Kaggle GPU: final offline notebook and heavier pretrained candidates where feasible.",
        "Cache candidate coordinates so fusion/refinement ablations remain cheap and reproducible.",
    ], 17, size=21)
    deck.bullet_slide("Immediate next steps and decision gates", [
        "1. Finish the data audit and rebuild reusable template artifacts on the main machine.",
        "2. Reproduce B0 from scratch and freeze its artifact bundle.",
        "3. Run and validate a Kaggle baseline notebook; late-submit its exact successful version.",
        "4. Add pretrained candidates and first measure oracle-pool TM.",
        "5. Proceed to fusion only if candidate sources are complementary.",
        "6. Require geometry v2 to eliminate the kink regression.",
    ], 18, size=20)
    deck.bullet_slide("Questions for the supervisor", [
        "Should the primary claim centre on segment fusion, with geometry v2 as projection/repair?",
        "Is 12-target temporal-safe CASP15 evaluation sufficient when paired with Kaggle private validation?",
        "Should the frozen final holdout be family-based, time-based, or target-based?",
        "Is a learned confidence gate necessary, or can a well-ablated heuristic gate suffice?",
        "Should success be framed as TM uplift, geometry validity, or an explicit two-axis claim?",
    ], 19, size=21)
    deck.bullet_slide("Appendix: evidence map", [
        "Pipeline results: reports/thesis_notes/results_summary.md",
        "Composite ablation: reports/thesis_notes/composite_ablation.md",
        "Refinement truthfulness: reports/thesis_notes/refine_ablation.md",
        "Top-1 reproduction and leakage: reproduce_top1.md + leakage_demo.md",
        "Full chronology: LOG.md; extension: PLAN.md + research_plan_review.md",
        "External papers and official Kaggle workflow: sources.md",
    ], 20, size=20)
    out = HERE / "supervisor_update.pptx"
    deck.save(out)
    print(f"wrote {out} ({len(deck.prs.slides)} slides)")


if __name__ == "__main__":
    main()
